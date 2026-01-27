"""
Unified document text extraction supporting PDF, DOCX, PPT, PPTX.
"""

from __future__ import annotations

import os
import re
from typing import Optional

import fitz

from app.utils.logger import get_logger

logger = get_logger(__name__)

# Supported file extensions
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".ppt", ".pptx"}


def get_supported_extensions() -> set:
    """Return set of supported file extensions."""
    return SUPPORTED_EXTENSIONS.copy()


def is_supported_file(filename: str) -> bool:
    """Check if the file extension is supported."""
    if not filename:
        return False
    ext = os.path.splitext(filename.lower())[1]
    return ext in SUPPORTED_EXTENSIONS


class DocumentExtractor:
    """Extract text from various document formats."""

    def extract_text(self, path: str) -> Optional[str]:
        """
        Extract text from a document based on its file extension.

        Supports: PDF, DOCX, DOC, PPT, PPTX

        Args:
            path: Path to the document file

        Returns:
            Extracted text or None if extraction fails
        """
        if not path or not os.path.exists(path):
            logger.info("document.not_found", path=path)
            return None

        ext = os.path.splitext(path.lower())[1]

        try:
            if ext == ".pdf":
                return self._extract_pdf(path)
            elif ext in {".docx", ".doc"}:
                return self._extract_docx(path)
            elif ext in {".ppt", ".pptx"}:
                return self._extract_pptx(path)
            else:
                logger.info("document.unsupported_format", path=path, extension=ext)
                return None
        except Exception as exc:
            logger.info("document.extraction_failed", path=path, error=str(exc))
            return None

    def _extract_pdf(self, path: str) -> Optional[str]:
        """Extract text from PDF using PyMuPDF."""
        try:
            doc = fitz.open(path)
        except Exception as exc:
            logger.info("pdf.open_failed", path=path, error=str(exc))
            return None

        if doc.is_encrypted:
            logger.info("pdf.encrypted", path=path)
            doc.close()
            return None

        text_parts = []
        for page in doc:
            text_parts.append(page.get_text("text"))
        doc.close()

        text = "\n".join(text_parts)
        return self._clean_text(text)

    def _extract_docx(self, path: str) -> Optional[str]:
        """Extract text from DOCX/DOC files."""
        try:
            from docx import Document
        except ImportError:
            logger.info("docx.library_not_installed")
            return None

        try:
            doc = Document(path)
            text_parts = []

            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))

            text = "\n".join(text_parts)
            return self._clean_text(text)
        except Exception as exc:
            logger.info("docx.extraction_failed", path=path, error=str(exc))
            return None

    def _extract_pptx(self, path: str) -> Optional[str]:
        """Extract text from PPT/PPTX files."""
        try:
            from pptx import Presentation
        except ImportError:
            logger.info("pptx.library_not_installed")
            return None

        try:
            prs = Presentation(path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []

                for shape in slide.shapes:
                    # Extract text from shapes with text frames
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))

                if slide_text:
                    text_parts.append(f"--- Slide {slide_num} ---")
                    text_parts.extend(slide_text)

            text = "\n".join(text_parts)
            return self._clean_text(text)
        except Exception as exc:
            logger.info("pptx.extraction_failed", path=path, error=str(exc))
            return None

    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text."""
        # Remove page number patterns
        cleaned = re.sub(r"Page\s+\d+", "", text, flags=re.IGNORECASE)
        # Reduce excessive newlines
        cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
        # Remove excessive whitespace
        cleaned = re.sub(r"[ \t]{3,}", "  ", cleaned)
        return cleaned.strip()
